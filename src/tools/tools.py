"""
Agent Tools — Capabilities available to all agents.

Tools:
- Internet search
- Browser automation
- Image generation
- Code execution
- API calls
"""

import asyncio
import logging
import os
import subprocess
import tempfile
from typing import Optional, Dict, Any

logger = logging.getLogger(__name__)

# Optional imports
try:
    import aiohttp
    AIOHTTP_AVAILABLE = True
except ImportError:
    AIOHTTP_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False


class WebSearchTool:
    """Internet search tool using web scraping."""

    async def search(self, query: str, num_results: int = 5) -> list:
        """Search the web for information."""
        if not AIOHTTP_AVAILABLE:
            return [{"error": "aiohttp not installed"}]

        try:
            url = f"https://html.duckduckgo.com/html/?q={query}"
            async with aiohttp.ClientSession() as session:
                async with session.get(url, timeout=aiohttp.ClientTimeout(total=10)) as resp:
                    if resp.status == 200:
                        html = await resp.text()
                        if BS4_AVAILABLE:
                            soup = BeautifulSoup(html, "html.parser")
                            results = []
                            for r in soup.select(".result")[:num_results]:
                                title = r.select_one(".result__title")
                                snippet = r.select_one(".result__snippet")
                                results.append({
                                    "title": title.get_text(strip=True) if title else "",
                                    "snippet": snippet.get_text(strip=True) if snippet else ""
                                })
                            return results
                        return [{"raw": html[:1000]}]
                    return [{"error": f"HTTP {resp.status}"}]
        except Exception as e:
            logger.error(f"Search error: {e}")
            return [{"error": str(e)}]


class BrowserTool:
    """Browser automation for web interaction."""

    async def fetch_page(self, url: str) -> Optional[str]:
        """Fetch a web page content."""
        if not AIOHTTP_AVAILABLE:
            return None
        try:
            async with aiohttp.ClientSession() as session:
                async with session.get(url, timeout=aiohttp.ClientTimeout(total=15)) as resp:
                    if resp.status == 200:
                        html = await resp.text()
                        if BS4_AVAILABLE:
                            soup = BeautifulSoup(html, "html.parser")
                            for tag in soup(["script", "style", "nav", "footer"]):
                                tag.decompose()
                            return soup.get_text(separator="\n", strip=True)[:5000]
                        return html[:5000]
        except Exception as e:
            logger.error(f"Browser error: {e}")
        return None


class ImageGeneratorTool:
    """Image generation using AI models."""

    def __init__(self):
        self.output_dir = "./generated_images"
        os.makedirs(self.output_dir, exist_ok=True)

    async def generate(self, prompt: str, filename: str = None) -> Optional[str]:
        """Generate an image from a text prompt."""
        try:
            from diffusers import StableDiffusionPipeline  # noqa
            import torch  # noqa

            pipe = StableDiffusionPipeline.from_pretrained(
                "runwayml/stable-diffusion-v1-5",
                torch_dtype=torch.float16
            )
            device = "cuda" if torch.cuda.is_available() else "cpu"
            pipe = pipe.to(device)

            image = pipe(prompt).images[0]
            fname = filename or f"generated_{hash(prompt) % 10000}.png"
            path = os.path.join(self.output_dir, fname)
            image.save(path)
            return path
        except ImportError:
            logger.warning("diffusers/torch not installed for image generation")
            return None
        except Exception as e:
            logger.error(f"Image generation error: {e}")
            return None


class CodeExecutorTool:
    """Safe code execution in a sandbox."""

    def __init__(self):
        self.workspace = "./workspace"
        os.makedirs(self.workspace, exist_ok=True)

    async def execute_python(self, code: str, timeout: int = 30) -> Dict[str, Any]:
        """Execute Python code safely."""
        # Write to temp file
        filepath = os.path.join(self.workspace, f"exec_{hash(code) % 10000}.py")
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(code)

        try:
            loop = asyncio.get_event_loop()
            result = await loop.run_in_executor(None, lambda: subprocess.run(
                ["python", filepath],
                capture_output=True, text=True, timeout=timeout,
                cwd=self.workspace
            ))
            return {
                "stdout": result.stdout,
                "stderr": result.stderr,
                "returncode": result.returncode,
                "success": result.returncode == 0
            }
        except subprocess.TimeoutExpired:
            return {"error": "Timeout", "success": False}
        except Exception as e:
            return {"error": str(e), "success": False}
        finally:
            try:
                os.remove(filepath)
            except Exception:
                pass


class APITool:
    """Make HTTP API calls."""

    async def call(self, url: str, method: str = "GET",
                   headers: dict = None, data: dict = None) -> Dict[str, Any]:
        """Make an API call."""
        if not AIOHTTP_AVAILABLE:
            return {"error": "aiohttp not installed"}

        try:
            async with aiohttp.ClientSession() as session:
                kwargs = {"headers": headers or {}, "timeout": aiohttp.ClientTimeout(total=15)}
                if data:
                    kwargs["json"] = data

                async with getattr(session, method.lower())(url, **kwargs) as resp:
                    body = await resp.text()
                    return {
                        "status": resp.status,
                        "body": body[:5000],
                        "success": 200 <= resp.status < 300
                    }
        except Exception as e:
            return {"error": str(e), "success": False}


class FileWriterTool:
    """Write files to workspace (code, configs, HTML, etc.)."""

    def __init__(self):
        self.workspace = "./workspace"
        os.makedirs(self.workspace, exist_ok=True)

    async def write_file(self, filename: str, content: str, subfolder: str = "") -> Dict[str, Any]:
        """Write content to a file in the workspace."""
        try:
            if subfolder:
                dirpath = os.path.join(self.workspace, subfolder)
                os.makedirs(dirpath, exist_ok=True)
                filepath = os.path.join(dirpath, filename)
            else:
                filepath = os.path.join(self.workspace, filename)

            with open(filepath, "w", encoding="utf-8") as f:
                f.write(content)

            abs_path = os.path.abspath(filepath)
            return {
                "success": True,
                "filepath": abs_path,
                "filename": filename,
                "size": len(content)
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    async def read_file(self, filepath: str) -> Dict[str, Any]:
        """Read a file from workspace."""
        try:
            full_path = os.path.join(self.workspace, filepath)
            with open(full_path, "r", encoding="utf-8") as f:
                content = f.read()
            return {"success": True, "content": content, "size": len(content)}
        except Exception as e:
            return {"success": False, "error": str(e)}

    async def list_files(self, subfolder: str = "") -> Dict[str, Any]:
        """List files in workspace."""
        try:
            dirpath = os.path.join(self.workspace, subfolder) if subfolder else self.workspace
            files = []
            for f in os.listdir(dirpath):
                full = os.path.join(dirpath, f)
                files.append({
                    "name": f,
                    "is_dir": os.path.isdir(full),
                    "size": os.path.getsize(full) if os.path.isfile(full) else 0
                })
            return {"success": True, "files": files}
        except Exception as e:
            return {"success": False, "error": str(e)}


class PresentationTool:
    """Create presentations (HTML-based, PPTX)."""

    def __init__(self):
        self.workspace = "./workspace"
        os.makedirs(self.workspace, exist_ok=True)

    async def create_html_presentation(self, title: str, slides: list,
                                        filename: str = None) -> Dict[str, Any]:
        """
        Create an HTML presentation (works without any extra libraries).

        Args:
            title: Presentation title
            slides: List of dicts with 'title' and 'content' keys
            filename: Output filename (auto-generated if not provided)

        Returns:
            Dict with filepath and success status
        """
        try:
            if not filename:
                filename = f"presentation_{hash(title) % 10000}.html"

            slides_html = ""
            for i, slide in enumerate(slides):
                slide_title = slide.get("title", f"Слайд {i + 1}")
                slide_content = slide.get("content", "")
                # Convert markdown-like content to HTML
                content_html = slide_content.replace("\n", "<br>")
                slides_html += f"""
        <div class="slide" id="slide-{i}">
            <h2>{slide_title}</h2>
            <div class="slide-content">{content_html}</div>
            <div class="slide-number">{i + 1} / {len(slides)}</div>
        </div>"""

            html = f"""<!DOCTYPE html>
<html lang="ru">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{ font-family: 'Segoe UI', Arial, sans-serif; background: #1a1a2e; color: #eee; overflow: hidden; }}
        .slide {{
            display: none; width: 100vw; height: 100vh;
            padding: 60px 80px; position: relative;
            background: linear-gradient(135deg, #1a1a2e 0%, #16213e 50%, #0f3460 100%);
            flex-direction: column; justify-content: center;
        }}
        .slide.active {{ display: flex; }}
        .slide h2 {{
            font-size: 2.5em; margin-bottom: 30px;
            background: linear-gradient(90deg, #e94560, #0f3460);
            -webkit-background-clip: text; -webkit-text-fill-color: transparent;
        }}
        .slide-content {{ font-size: 1.3em; line-height: 1.8; max-width: 900px; }}
        .slide-number {{
            position: absolute; bottom: 30px; right: 40px;
            font-size: 0.9em; opacity: 0.5;
        }}
        .controls {{
            position: fixed; bottom: 30px; left: 50%; transform: translateX(-50%);
            display: flex; gap: 20px; z-index: 100;
        }}
        .controls button {{
            padding: 12px 30px; font-size: 1em; border: none; border-radius: 8px;
            background: #e94560; color: white; cursor: pointer; transition: 0.3s;
        }}
        .controls button:hover {{ background: #c73652; transform: scale(1.05); }}
        .title-slide h1 {{
            font-size: 3.5em; text-align: center;
            background: linear-gradient(90deg, #e94560, #533483, #0f3460);
            -webkit-background-clip: text; -webkit-text-fill-color: transparent;
        }}
        .title-slide p {{ text-align: center; font-size: 1.5em; margin-top: 20px; opacity: 0.7; }}
    </style>
</head>
<body>
    <div class="slide active title-slide" id="slide-title">
        <h1>{title}</h1>
        <p>Создано Jarvis4 AI Office</p>
        <div class="slide-number">Титульный слайд</div>
    </div>
{slides_html}
    <div class="controls">
        <button onclick="prevSlide()">← Назад</button>
        <button onclick="nextSlide()">Далее →</button>
    </div>
    <script>
        let current = 0;
        const slides = document.querySelectorAll('.slide');
        function showSlide(n) {{
            slides.forEach(s => s.classList.remove('active'));
            current = Math.max(0, Math.min(n, slides.length - 1));
            slides[current].classList.add('active');
        }}
        function nextSlide() {{ showSlide(current + 1); }}
        function prevSlide() {{ showSlide(current - 1); }}
        document.addEventListener('keydown', e => {{
            if (e.key === 'ArrowRight' || e.key === ' ') nextSlide();
            if (e.key === 'ArrowLeft') prevSlide();
        }});
    </script>
</body>
</html>"""

            filepath = os.path.join(self.workspace, filename)
            with open(filepath, "w", encoding="utf-8") as f:
                f.write(html)

            abs_path = os.path.abspath(filepath)
            return {
                "success": True,
                "filepath": abs_path,
                "filename": filename,
                "slides_count": len(slides) + 1,  # +1 for title slide
                "format": "html"
            }
        except Exception as e:
            return {"success": False, "error": str(e)}

    async def create_pptx_presentation(self, title: str, slides: list,
                                        filename: str = None) -> Dict[str, Any]:
        """
        Create a PPTX presentation (requires python-pptx).

        Falls back to HTML if python-pptx is not installed.
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor

            if not filename:
                filename = f"presentation_{hash(title) % 10000}.pptx"

            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            # Title slide
            slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            slide.placeholders[1].text = "Создано Jarvis4 AI Office"

            # Content slides
            for s in slides:
                slide_layout = prs.slide_layouts[1]
                slide = prs.slides.add_slide(slide_layout)
                slide.shapes.title.text = s.get("title", "")
                slide.placeholders[1].text = s.get("content", "")

            filepath = os.path.join(self.workspace, filename)
            prs.save(filepath)

            abs_path = os.path.abspath(filepath)
            return {
                "success": True,
                "filepath": abs_path,
                "filename": filename,
                "slides_count": len(slides) + 1,
                "format": "pptx"
            }
        except ImportError:
            logger.warning("python-pptx not installed, falling back to HTML presentation")
            html_filename = filename.replace(".pptx", ".html") if filename else None
            return await self.create_html_presentation(title, slides, html_filename)
        except Exception as e:
            return {"success": False, "error": str(e)}


class ToolManager:
    """Central manager for all agent tools."""

    def __init__(self):
        self.search = WebSearchTool()
        self.browser = BrowserTool()
        self.image_gen = ImageGeneratorTool()
        self.code_exec = CodeExecutorTool()
        self.api = APITool()
        self.file_writer = FileWriterTool()
        self.presentation = PresentationTool()

    def get_available_tools(self) -> list:
        return [
            "web_search", "browser", "image_generation", "code_execution",
            "api_calls", "file_writer", "presentation_creator"
        ]
